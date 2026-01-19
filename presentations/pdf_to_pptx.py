"""
PDF to PowerPoint Converter

Converts NotebookLM PDF slides into an editable PowerPoint presentation.
Each PDF page becomes a slide background image, allowing you to add
editable text boxes, notes, and additional slides.

Requirements:
    pip install pdf2image python-pptx Pillow

On Windows, you also need poppler:
    1. Download from: https://github.com/osber/pdf2image/releases
    2. Extract to C:\poppler
    3. Add C:\poppler\bin to your PATH

Or use conda:
    conda install -c conda-forge poppler

Usage:
    python pdf_to_pptx.py input.pdf output.pptx
"""

import sys
import os
from pathlib import Path

try:
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from PIL import Image
    import tempfile
except ImportError as e:
    print(f"Missing required package: {e}")
    print("\nInstall requirements with:")
    print("    pip install pdf2image python-pptx Pillow")
    print("\nOn Windows, also install poppler:")
    print("    conda install -c conda-forge poppler")
    print("    OR download from https://github.com/osber/pdf2image/releases")
    sys.exit(1)


def pdf_to_pptx(pdf_path: str, output_path: str, dpi: int = 200):
    """
    Convert a PDF file to an editable PowerPoint presentation.

    Args:
        pdf_path: Path to input PDF file
        output_path: Path for output PPTX file
        dpi: Resolution for image extraction (higher = better quality, larger file)
    """

    print(f"Converting: {pdf_path}")
    print(f"Output: {output_path}")
    print(f"DPI: {dpi}")
    print()

    # Convert PDF pages to images
    print("Extracting PDF pages as images...")
    try:
        images = convert_from_path(pdf_path, dpi=dpi)
    except Exception as e:
        print(f"\nError converting PDF: {e}")
        print("\nIf you're on Windows, make sure poppler is installed:")
        print("    conda install -c conda-forge poppler")
        print("    OR download from https://github.com/osber/pdf2image/releases")
        sys.exit(1)

    print(f"Found {len(images)} pages")

    # Create PowerPoint presentation
    # Use widescreen 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # Blank slide layout
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Create a temporary directory for images
    with tempfile.TemporaryDirectory() as temp_dir:
        for i, image in enumerate(images):
            print(f"Processing slide {i + 1}/{len(images)}...")

            # Save image to temp file
            img_path = os.path.join(temp_dir, f"slide_{i + 1}.png")
            image.save(img_path, "PNG")

            # Add slide
            slide = prs.slides.add_slide(blank_layout)

            # Add image as background (full slide size)
            slide.shapes.add_picture(
                img_path,
                Inches(0),
                Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )

    # Save PowerPoint
    print(f"\nSaving PowerPoint to: {output_path}")
    prs.save(output_path)
    print("Done!")

    # Print instructions
    print("\n" + "=" * 60)
    print("SUCCESS! Your editable PowerPoint is ready.")
    print("=" * 60)
    print("\nTo edit in PowerPoint:")
    print("1. Open the .pptx file")
    print("2. The images are backgrounds - you can add text boxes on top")
    print("3. Add new slides at the end for 'Decisions Needed' section")
    print("4. Add speaker notes as needed")
    print("\nTip: To edit text in the images, add a text box and")
    print("     position it over the area you want to 'replace'")


def extract_images_only(pdf_path: str, output_dir: str, dpi: int = 200):
    """
    Extract PDF pages as individual PNG images.

    Args:
        pdf_path: Path to input PDF file
        output_dir: Directory to save images
        dpi: Resolution for image extraction
    """

    print(f"Extracting images from: {pdf_path}")
    print(f"Output directory: {output_dir}")
    print(f"DPI: {dpi}")
    print()

    # Create output directory
    Path(output_dir).mkdir(parents=True, exist_ok=True)

    # Convert PDF pages to images
    print("Extracting PDF pages...")
    images = convert_from_path(pdf_path, dpi=dpi)
    print(f"Found {len(images)} pages")

    # Save each image
    for i, image in enumerate(images):
        output_path = os.path.join(output_dir, f"slide_{i + 1:02d}.png")
        print(f"Saving: {output_path}")
        image.save(output_path, "PNG")

    print("\nDone! Images saved to:", output_dir)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage:")
        print("  Convert PDF to PowerPoint:")
        print("    python pdf_to_pptx.py input.pdf output.pptx")
        print()
        print("  Extract images only:")
        print("    python pdf_to_pptx.py input.pdf output_folder/ --images-only")
        print()
        print("  Custom DPI (default 200):")
        print("    python pdf_to_pptx.py input.pdf output.pptx --dpi 300")
        sys.exit(1)

    pdf_path = sys.argv[1]
    output_path = sys.argv[2]

    # Check for flags
    images_only = "--images-only" in sys.argv

    # Get DPI
    dpi = 200
    if "--dpi" in sys.argv:
        dpi_index = sys.argv.index("--dpi") + 1
        if dpi_index < len(sys.argv):
            dpi = int(sys.argv[dpi_index])

    # Validate input
    if not os.path.exists(pdf_path):
        print(f"Error: File not found: {pdf_path}")
        sys.exit(1)

    # Run conversion
    if images_only:
        extract_images_only(pdf_path, output_path, dpi)
    else:
        pdf_to_pptx(pdf_path, output_path, dpi)
